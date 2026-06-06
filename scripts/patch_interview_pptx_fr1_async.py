#!/usr/bin/env python3
"""Layout + async evidence for FR1 slide (slide 3) in the pro deck."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
DECK = ROOT / "docs" / "deliverables" / "v0.1" / "zooplus-assistant-interview-15min-pro.pptx"

INK = RGBColor(0x33, 0x41, 0x55)
TITLE_INK = RGBColor(0x0A, 0x25, 0x40)
PALE = RGBColor(0xBF, 0xDB, 0xFE)
CODE_BG = RGBColor(0x1E, 0x29, 0x3B)
CODE_FG = RGBColor(0xE2, 0xE8, 0xF0)
ACCENT = RGBColor(0x25, 0x63, 0xEB)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
JSON_BG = RGBColor(0xE0, 0xF2, 0xFE)

LEFT_BULLETS = [
    ("FR1 — Async FastAPI · POST /chat { site_id, query }", True),
    ("async def chat() → await handle_chat()  (routes/chat.py)", False),
    ("Orchestrator async; Chroma/intent via asyncio.to_thread", False),
    ("Deploy: uvicorn ASGI  ·  Test B1: iscoroutinefunction(chat)", False),
]

CODE_SNIPPET = (
    '@router.post("/chat")\n'
    "async def chat(body: ChatRequest):\n"
    "    return await handle_chat(body)\n"
    "\n"
    "# B1  tests/acceptance/test_coding_task_brief.py\n"
    "assert inspect.iscoroutinefunction(chat)"
)

JSON_BODY = (
    '{\n'
    '  "site_id": 3,\n'
    '  "query": "What\'s the best dry food for a puppy with a sensitive stomach?"\n'
    "}"
)

REMOVE_NAMES = {
    "Panel AsyncEvidence",
    "TextBox AsyncEvidence",
    "TextBox B1Badge",
}


def _set_bullets(text_frame, lines: list[tuple[str, bool]], *, size: int = 15) -> None:
    text_frame.clear()
    for i, (line, bold) in enumerate(lines):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = line
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = TITLE_INK if bold else INK
        p.space_after = Pt(6)


def _set_mono(text_frame, text: str, *, size: int = 11, color: RGBColor = INK) -> None:
    text_frame.clear()
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = line
        p.font.name = "Consolas"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(1)


def _rgb(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _remove_shapes(slide, names: set[str]) -> None:
    to_drop = [sh for sh in slide.shapes if sh.name in names]
    for sh in to_drop:
        sp = sh._element
        sp.getparent().remove(sp)


def _find_slide(prs: Presentation, title_prefix: str) -> int:
    for i, slide in enumerate(prs.slides):
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip().startswith(title_prefix):
                return i
    raise SystemExit(f"Slide not found: {title_prefix!r}")


def _get_or_add_code_panel(slide):
    for sh in slide.shapes:
        if sh.name == "Panel AsyncEvidenceV2":
            return sh
    panel = slide.shapes.add_shape(
        1, Inches(6.62), Inches(4.72), Inches(5.58), Inches(2.05)
    )
    panel.name = "Panel AsyncEvidenceV2"
    _rgb(panel, CODE_BG)
    box = slide.shapes.add_textbox(Inches(6.74), Inches(4.82), Inches(5.35), Inches(1.85))
    box.name = "TextBox AsyncEvidenceV2"
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    h = tf.paragraphs[0]
    h.text = "Async evidence (source) — B1 PASS"
    h.font.name = "Arial"
    h.font.size = Pt(11)
    h.font.bold = True
    h.font.color.rgb = GREEN
    body = tf.add_paragraph()
    body.text = CODE_SNIPPET
    body.font.name = "Consolas"
    body.font.size = Pt(10)
    body.font.color.rgb = CODE_FG
    return panel


def fix_fr1_slide(prs: Presentation) -> None:
    slide = prs.slides[_find_slide(prs, "API endpoint")]
    _remove_shapes(slide, REMOVE_NAMES | {"Panel AsyncEvidenceV2", "TextBox AsyncEvidenceV2"})

    # --- Left column zones (no overlap) ---
    # Panel bg
    for sh in slide.shapes:
        if sh.name == "Rectangle 11":
            sh.left, sh.top = Inches(0.64), Inches(1.38)
            sh.width, sh.height = Inches(5.35), Inches(5.55)
            _rgb(sh, RGBColor(0xF8, 0xFA, 0xFC))

    # Bullets — top third
    for sh in slide.shapes:
        if sh.name == "TextBox 12" and sh.has_text_frame:
            sh.left, sh.top = Inches(0.78), Inches(1.52)
            sh.width, sh.height = Inches(5.05), Inches(1.95)
            _set_bullets(sh.text_frame, LEFT_BULLETS, size=14)
        if sh.name == "TextBox 9" and sh.has_text_frame:
            sh.text_frame.paragraphs[0].text = (
                "FR1 — Functional Requirement 1: async FastAPI + POST /chat"
            )

    # Mandatory JSON — middle band
    for sh in slide.shapes:
        if sh.name == "Rectangle 13":
            sh.left, sh.top = Inches(0.78), Inches(3.62)
            sh.width, sh.height = Inches(5.05), Inches(1.35)
            _rgb(sh, JSON_BG)
        if sh.name == "TextBox 14" and sh.has_text_frame:
            sh.left, sh.top = Inches(0.92), Inches(3.78)
            sh.width, sh.height = Inches(4.75), Inches(1.05)
            _set_mono(sh.text_frame, JSON_BODY, size=12, color=TITLE_INK)

    # Swagger CTA — bottom left
    cta_names = {s.name for s in slide.shapes}
    if "TextBox SwaggerCTA" not in cta_names:
        cta = slide.shapes.add_textbox(Inches(0.78), Inches(5.15), Inches(5.05), Inches(0.55))
        cta.name = "TextBox SwaggerCTA"
    else:
        cta = next(s for s in slide.shapes if s.name == "TextBox SwaggerCTA")
    cta.left, cta.top = Inches(0.78), Inches(5.15)
    cta.width, cta.height = Inches(5.05), Inches(0.55)
    p = cta.text_frame.paragraphs[0]
    p.text = "Live proof: /docs — Try it out with the mandatory brief query above"
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    # --- Right column: Swagger top, code bottom ---
    for sh in slide.shapes:
        if sh.name == "Picture 16":
            sh.left, sh.top = Inches(6.62), Inches(1.45)
            sh.width, sh.height = Inches(5.58), Inches(3.15)
        if sh.name == "Rounded Rectangle 15":
            sh.left, sh.top = Inches(6.55), Inches(1.38)
            sh.width, sh.height = Inches(5.72), Inches(5.55)
            _rgb(sh, RGBColor(0xFF, 0xFF, 0xFF))

    _get_or_add_code_panel(slide)


def main() -> None:
    if not DECK.is_file():
        raise SystemExit(f"Missing deck: {DECK}")
    prs = Presentation(str(DECK))
    fix_fr1_slide(prs)
    prs.save(str(DECK))
    print(f"Fixed slide 3 layout in {DECK}")


if __name__ == "__main__":
    main()
