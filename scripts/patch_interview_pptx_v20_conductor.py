#!/usr/bin/env python3
"""Update slides 5–9 diagrams (v2.1.6) and release-progress closing slide."""

from __future__ import annotations

import shutil
from pathlib import Path

from patch_interview_pptx_rag_slides import (  # type: ignore[import-not-found]
    PANEL_BG,
    draw_pro_chrome,
    renumber_footer,
)
from patch_interview_pptx_rag_slides import _bullets as _rag_bullets  # type: ignore[import-not-found]
from patch_interview_pptx_rag_slides import _rgb as _rgb_panel  # type: ignore[import-not-found]
from patch_interview_pptx_v14_live_loop import (  # type: ignore[import-not-found]
    DECK,
    _fill,
    _find_slide,
    _fix_left_column,
    _right_column_bullets,
)
from pptx import Presentation
from pptx.util import Inches

SLIDE5_RIGHT_V20 = [
    ("Query path (hybrid default)", True),
    ("1. Hard filter site_id in Chroma metadata", False),
    ("2. Pull candidates — pool scales with price band / requested count", False),
    ("3. Score BM25 + business signals on candidate pool", False),
    ("4. Rerank → default 4 picks; shopper can ask up to 20", False),
    ("Stream delivery (v2.1.6)", True),
    ("• product_batch NDJSON — UI reveals product cards in chunks of 4", False),
    ("• A/B vector-only: ZOOPLUS_RETRIEVAL_MODE=vector", False),
]

SLIDE6_LEFT_V20 = [
    ("Ingest (offline)", True),
    ("• Source JSON → data/raw (300 variants) · strip HTML.", False),
    ("• Chroma per article_id · routing_lexicon.json at ingest.", False),
    ("Online path (catalog lane)", True),
    ("• probe_instant_lane: phrase index + playbook → social | catalog.", False),
    ("• Hybrid retrieve → EUR price filter → resolve count from query.", False),
    ("NDJSON stream (/chat/stream)", True),
    ("• typing → chunk* → product_batch* → done (+ meta).", False),
    ("• Social/help: no catalog progress; dedupe final vs live chunks.", False),
]

SLIDE6_RIGHT_V20 = [
    ("Grounding & answer", True),
    ("• retrieved_products[] from same site_id hits only.", False),
    ("• Off-topic → empty list + polite decline (FR4).", False),
    ("• Synthesis: per-agent OpenCode LLM — never invent SKUs.", False),
    ("v2.1.6 — playbook + smart limits", True),
    ("• 4 picks default · shopper can ask for more (e.g. 10 opciones, cap 20).", False),
    ("• product_batch stream — cards arrive in chunks of 4 in the UI.", False),
    ("• social_phrases.yaml index + playbook auto-learn (help/greeting).", False),
    ("• Dynamic species inference (iguana, capibara…) — not a fixed list.", False),
]

SLIDE7_FLOW_V20 = (
    "Request flow: UI → probe (phrase index) → intent-agent → "
    "social-agent | catalog prefetch → chunks + product_batch → done"
)

SLIDE7_LEFT_V20 = [
    ("Invisible conductor (releases v2.0+)", True),
    ("• Internal MD playbook — conductor maintains silently.", False),
    ("• Social turn → one bubble; help asks never hit catalog progress.", False),
    ("• Catalog turn → ack from query (species, price, lang, count).", False),
    ("• Dedupe: no re-greeting / no repeated scope in final answer.", False),
    ("Stream + fast intent (v2.1.3+)", True),
    ("• Policy + intent-agent first; conductor opt-in (lower latency).", False),
    ("• Parallel social/catalog · ZOOPLUS_STREAM_MODE=conductor.", False),
]

SLIDE8_LEFT_V20 = [
    ("OpenCode specialists", True),
    ("• zooplus-conductor → playbook + stream orchestration (invisible).", False),
    ("• social-agent → shopper voice; mid-session help without re-intro.", False),
    ("• intent / topic-guard → scope filter.", False),
    ("• rag + logic + synthesis → grounded catalog answer.", False),
    ("Phrase index + language", True),
    ("• ~90 curated social phrases (ES/EN/DE/FR) — fast in-memory match.", False),
    ("• Agent mirrors shopper language; playbook learns novel phrases.", False),
]

SLIDE9_LEFT_V20 = [
    ("FR4 guardrails + live demo", True),
    ("• Pet catalog only · default-deny off-topic.", False),
    ("Demo — v2.1.6 smart loop", True),
    ("• A: me puedes ayudar → social help, no catalog progress.", False),
    ("• B: y para iguanas → scope reply, no duplicate Hola intro.", False),
    ("• C: dame 10 opciones perros → up to 10 cards in batches.", False),
    ("• http://127.0.0.1:8090/ui/ — shop Spain (15).", False),
]

PROGRESS_LEFT = [
    ("Foundation — v0.1.0 → v1.0", True),
    ("v0.1.0 — Coding Task PoC: FR1–5, hybrid RAG, FR4 guardrails.", False),
    ("v0.1.3 — Conductor-first routing; /chat/stream; per-agent OpenCode.", False),
    ("v1.0 — Stable tag: catalog lexicon at ingest, cache, multilingual.", False),
]

PROGRESS_RIGHT = [
    ("Shopper experience — v1.4 → v2.1.6 (now)", True),
    ("v2.1 — Playbook MD; social/catalog probe; agent-multilingual.", False),
    ("v2.1.3 — Fast intent-first stream; conductor intent opt-in.", False),
    ("v2.1.4–5 — Species inference; help/saludo detect + dedupe.", False),
    ("v2.1.6 — Dynamic count (4→20); phrase index; product_batch UI.", False),
]


def _find_progress_slide(prs: Presentation) -> int | None:
    for i, slide in enumerate(prs.slides):
        for sh in slide.shapes:
            if sh.has_text_frame and (sh.text or "").strip().startswith("Release progress"):
                return i
    return None


def _append_blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


def _set_progress_footer(slide) -> None:
    footer = "zooplus Assistant — releases v2.1.6 · progress since v0.1.0"
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        text = (sh.text or "").strip()
        if "zooplus Assistant" in text and sh.top and sh.top > Inches(6.8):
            sh.text_frame.paragraphs[0].text = footer


def build_release_progress_slide(slide, *, page: int) -> None:
    draw_pro_chrome(
        slide,
        title="Release progress — v0.1.0 → v2.1.6",
        subtitle="Same five FRs · smarter social routing, flexible picks, chunked UX",
        badge="Progress · releases",
        page=page,
    )
    panel = slide.shapes.add_shape(1, Inches(0.64), Inches(1.38), Inches(12.0), Inches(5.45))
    _rgb_panel(panel, PANEL_BG)
    _rag_bullets(
        slide,
        Inches(0.82),
        Inches(1.55),
        Inches(5.85),
        Inches(5.1),
        PROGRESS_LEFT,
        size=15,
    )
    _rag_bullets(
        slide,
        Inches(6.7),
        Inches(1.55),
        Inches(5.75),
        Inches(5.1),
        PROGRESS_RIGHT,
        size=15,
    )
    _set_progress_footer(slide)


def _remove_stale_rag_middle_column(slide) -> None:
    """Drop legacy centre column on slide 6 (v1.4 status* stream text)."""
    for sh in list(slide.shapes):
        if not sh.has_text_frame:
            continue
        text = sh.text or ""
        if sh.left and 700_000 < sh.left < 2_000_000 and "status* → products" in text:
            el = sh._element
            el.getparent().remove(el)


def ensure_release_progress_slide(prs: Presentation) -> None:
    idx = _find_progress_slide(prs)
    if idx is None:
        slide = _append_blank_slide(prs)
        build_release_progress_slide(slide, page=len(prs.slides))
    else:
        slide = prs.slides[idx]
        for sh in list(slide.shapes):
            el = sh._element
            el.getparent().remove(el)
        build_release_progress_slide(slide, page=idx + 1)
    renumber_footer(prs)


def main() -> None:
    if not DECK.is_file():
        raise SystemExit(f"Missing deck: {DECK}")
    prs = Presentation(str(DECK))

    i5 = _find_slide(prs, "Why hybrid")
    right5 = _right_column_bullets(prs.slides[i5])
    if right5 is not None:
        _fill(right5, SLIDE5_RIGHT_V20, size=14)

    i6 = _find_slide(prs, "RAG strategy")
    _remove_stale_rag_middle_column(prs.slides[i6])
    _fix_left_column(prs.slides[i6], SLIDE6_LEFT_V20, size=14)
    right6 = _right_column_bullets(prs.slides[i6])
    if right6 is not None:
        _fill(right6, SLIDE6_RIGHT_V20, size=14)

    i7 = _find_slide(prs, "Agentic architecture", "Agentic routing", "Live loop")
    _fix_left_column(prs.slides[i7], SLIDE7_LEFT_V20, size=14)
    for sh in prs.slides[i7].shapes:
        if not sh.has_text_frame:
            continue
        text = sh.text or ""
        if "How each request" in text:
            sh.text_frame.paragraphs[0].text = SLIDE7_FLOW_V20
        elif "Conductor-first" in text or "catalog lexicon · stream status" in text:
            sh.text_frame.paragraphs[0].text = (
                "v2.1.6 · phrase-index probe · fast intent · product_batch stream"
            )

    i8 = _find_slide(prs, "Agents")
    _fix_left_column(prs.slides[i8], SLIDE8_LEFT_V20, size=13)

    i9 = _find_slide(prs, "Guardrails", "FR4 guardrails")
    _fix_left_column(prs.slides[i9], SLIDE9_LEFT_V20, size=14)

    ensure_release_progress_slide(prs)

    tmp = DECK.with_name(f"{DECK.stem}_v20{DECK.suffix}")
    prs.save(str(tmp))
    try:
        shutil.copy2(tmp, DECK)
        tmp.unlink(missing_ok=True)
        print(f"Updated v2.1.6 conductor + release progress slide in {DECK} ({len(prs.slides)} slides)")
    except PermissionError:
        print(f"PPT is open — saved to {tmp}")


if __name__ == "__main__":
    main()
