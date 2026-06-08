#!/usr/bin/env python3
"""Insert two RAG/hybrid slides into the pro interview deck (after slide 4)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
DECK = ROOT / "docs" / "deliverables" / "v0.1" / "zooplus-assistant-interview-15min-pro.pptx"

NAVY = RGBColor(0x0A, 0x25, 0x40)
BLUE = RGBColor(0x25, 0x63, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALE = RGBColor(0xBF, 0xDB, 0xFE)
INK = RGBColor(0x33, 0x41, 0x55)
TITLE_INK = RGBColor(0x0A, 0x25, 0x40)
PANEL_BG = RGBColor(0xF8, 0xFA, 0xFC)


def _rgb(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _text(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = 20,
    bold: bool = False,
    color: RGBColor = INK,
    align=PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = "Arial"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color


def _bullets(
    slide,
    left,
    top,
    width,
    height,
    lines: list[tuple[str, bool]],
    *,
    size: int = 18,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, (line, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = TITLE_INK if bold else INK
        p.space_after = Pt(6)


def draw_pro_chrome(slide, *, title: str, subtitle: str, badge: str, page: int) -> None:
    sw, sh = Inches(13.333), Inches(7.5)
    _rgb(slide.shapes.add_shape(1, Inches(0.47), 0, sw - Inches(0.47), sh), WHITE)
    _rgb(slide.shapes.add_shape(1, 0, 0, Inches(0.42), sh), NAVY)
    _rgb(slide.shapes.add_shape(1, Inches(0.42), 0, Inches(0.05), sh), BLUE)
    _rgb(slide.shapes.add_shape(1, Inches(0.47), 0, sw - Inches(0.47), Inches(1.22)), NAVY)
    _rgb(slide.shapes.add_shape(1, Inches(0.47), Inches(1.22), sw - Inches(0.47), Inches(0.06)), BLUE)
    _rgb(slide.shapes.add_shape(1, Inches(0.47), Inches(7.02), sw - Inches(0.47), Inches(0.48)), NAVY)
    _rgb(slide.shapes.add_shape(1, 0, 0, Inches(0.42), sh), NAVY)
    _rgb(slide.shapes.add_shape(1, Inches(0.42), 0, Inches(0.05), sh), BLUE)

    _text(slide, Inches(0.72), Inches(0.2), Inches(10), Inches(0.62), title, size=34, bold=True, color=WHITE)
    _text(slide, Inches(0.72), Inches(0.78), Inches(8.8), Inches(0.38), subtitle, size=17, color=PALE)
    badge_shape = slide.shapes.add_shape(1, Inches(11.2), Inches(0.32), Inches(1.75), Inches(0.48))
    _rgb(badge_shape, BLUE)
    _text(
        slide,
        Inches(11.28),
        Inches(0.38),
        Inches(1.6),
        Inches(0.4),
        badge,
        size=13,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    _text(
        slide,
        Inches(0.72),
        Inches(7.1),
        Inches(10.2),
        Inches(0.36),
        "zooplus Assistant — Coding Task PoC v0.1",
        size=12,
        color=PALE,
    )
    _text(slide, Inches(12.2), Inches(7.08), Inches(0.95), Inches(0.36), str(page), size=15, bold=True, color=WHITE)


def insert_slide_at(prs: Presentation, index: int):
    layout = prs.slide_layouts[6]
    prs.slides.add_slide(layout)
    xml = prs.slides._sldIdLst
    elems = list(xml)
    new_elem = elems[-1]
    xml.remove(new_elem)
    xml.insert(index, new_elem)
    return prs.slides[index]


def renumber_footer(prs: Presentation) -> None:
    for i, slide in enumerate(prs.slides, start=1):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            if sh.top and sh.top > Inches(6.9) and sh.width and sh.width < Inches(1.2):
                t = sh.text_frame.text.strip()
                if t.isdigit():
                    sh.text_frame.paragraphs[0].text = str(i)


def build_hybrid_slide(slide) -> None:
    draw_pro_chrome(
        slide,
        title="Why hybrid retrieval?",
        subtitle="Chroma vectors + BM25 lexical + business rerank — not vector-only",
        badge="RAG · Hybrid",
        page=5,
    )
    panel = slide.shapes.add_shape(1, Inches(0.64), Inches(1.38), Inches(12.0), Inches(5.45))
    _rgb(panel, PANEL_BG)
    _bullets(
        slide,
        Inches(0.82),
        Inches(1.55),
        Inches(5.6),
        Inches(5.1),
        [
            ("Problem we saw in the catalog", True),
            ("• Vector search alone misses exact tokens (brands, SKUs, “grain-free”).", False),
            ("• Pure keyword search misses intent (“puppy sensitive stomach”).", False),
            ("Why this hybrid design", True),
            ("• Chroma: semantic candidates from embedded product text (local, no API key).", False),
            ("• BM25 on the same candidate pool: lexical match without a second DB.", False),
            ("• Fuse 50% vector + 35% BM25 + 15% rating/sales/stock.", False),
            ("PoC trade-off", True),
            ("• Zero extra infra vs hosted vector DB; A/B via ZOOPLUS_RETRIEVAL_MODE=vector.", False),
        ],
        size=17,
    )
    _bullets(
        slide,
        Inches(6.55),
        Inches(1.55),
        Inches(5.9),
        Inches(5.1),
        [
            ("Query path (hybrid default)", True),
            ("1. Hard filter site_id in Chroma metadata", False),
            ("2. Pull ≥24 vector candidates", False),
            ("3. Score BM25 on candidate documents", False),
            ("4. Rerank → default 4; shopper can ask more (cap 20)", False),
            ("5. Stream: product_batch chunks of 4 in UI (v2.1.6)", False),
            ("When vector-only is enough", True),
            ("• Short brand lookups or regression tests", False),
            ("• Hybrid wins on natural-language shopper queries", False),
        ],
        size=17,
    )


def build_rag_pipeline_slide(slide) -> None:
    draw_pro_chrome(
        slide,
        title="RAG strategy — end to end",
        subtitle="Ingest → retrieve → ground → synthesize (catalog-only, site-scoped)",
        badge="FR3 · RAG",
        page=6,
    )
    panel = slide.shapes.add_shape(1, Inches(0.64), Inches(1.38), Inches(12.0), Inches(5.45))
    _rgb(panel, PANEL_BG)
    _bullets(
        slide,
        Inches(0.82),
        Inches(1.55),
        Inches(5.85),
        Inches(5.1),
        [
            ("Ingest (offline)", True),
            ("• Source of truth: instructions JSON → data/raw (300 variants).", False),
            ("• Strip HTML, one Chroma doc per article_id row.", False),
            ("• Metadata: site_id, pet_type, price, stock, ingredients flag.", False),
            ("• CLI: python -m cli ingest → artifacts/index/chroma", False),
            ("Online retrieval", True),
            ("• Topic guard ALLOW → process lane search_catalog().", False),
            ("• Optional price-band filter parsed from query (EUR range).", False),
            ("• Default 4 picks; resolve_recommendation_count() up to 20.", False),
        ],
        size=16,
    )
    _bullets(
        slide,
        Inches(6.7),
        Inches(1.55),
        Inches(5.75),
        Inches(5.1),
        [
            ("Grounding & answer", True),
            ("• retrieved_products[] always from same site_id hits.", False),
            ("• Off-topic → empty list + polite decline (FR4).", False),
            ("• Synthesis: OpenCode agents or template — never invent SKUs.", False),
            ("• Stream: typing → chunk* → product_batch* → done.", False),
            ("Production path", True),
            ("• Swap Chroma for managed vector DB; keep hybrid fusion logic.", False),
            ("• Scheduled re-ingest when catalog JSON changes (roadmap P1).", False),
        ],
        size=16,
    )


def main() -> None:
    if not DECK.is_file():
        raise SystemExit(f"Missing deck: {DECK}")
    prs = Presentation(str(DECK))
    if len(prs.slides) >= 13:
        print(f"Deck already has {len(prs.slides)} slides — skip insert.")
        return

    s5 = insert_slide_at(prs, 4)
    build_hybrid_slide(s5)
    s6 = insert_slide_at(prs, 5)
    build_rag_pipeline_slide(s6)
    renumber_footer(prs)
    prs.save(str(DECK))
    print(f"Updated {DECK} — now {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
