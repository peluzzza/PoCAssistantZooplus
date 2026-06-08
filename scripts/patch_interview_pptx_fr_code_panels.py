#!/usr/bin/env python3
"""Add FR evidence code panels to slides 4, 6, 9, 10 (slide 3 = fr1_async)."""

from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
DECK = ROOT / "docs" / "deliverables" / "v0.1" / "zooplus-assistant-interview-15min-pro.pptx"

CODE_BG = RGBColor(0x1E, 0x29, 0x3B)
CODE_FG = RGBColor(0xE2, 0xE8, 0xF0)
GREEN = RGBColor(0x16, 0xA3, 0x4A)

FR24_CODE = (
    "# models/chat.py — FR2 response contract\n"
    "class ChatResponse(BaseModel):\n"
    "    answer: str\n"
    "    retrieved_products: list[RetrievedProduct]\n"
    "\n"
    "# lanes/process.py — FR3 catalog-only RAG\n"
    "cap = resolve_recommendation_count(query)\n"
    "hits = search_catalog(query, site_id, n_results=pool_n)\n"
    "products = [_to_retrieved_product(h) for h in hits][:cap]"
)

FR3_RAG_CODE = (
    "# rag/hybrid.py — site_id filter + fusion\n"
    "hits = chroma_query(index_dir(), query, site_id, n_results=pool)\n"
    "hit['hybrid_score'] = fuse_hybrid_scores(vec, bm25, business)\n"
    "return hits[:n_results]\n"
    "\n"
    "# Source of truth: product_catalog_dataset.json only\n"
    "python -m cli ingest  →  artifacts/index/chroma"
)

FR4_CODE = (
    "# guardian/engine.py — FR4 default-deny\n"
    "def topic_check(query, site_id=3) -> TopicDecision:\n"
    "    if OFF_TOPIC_PATTERN.search(query):\n"
    "        return TopicDecision(decision='DECLINE', ...)\n"
    "\n"
    "# constraints.yaml\n"
    "topic_guard_mode: default_deny\n"
    "decline_intents: [weather, non_pet_products, ...]"
)

FR5_CODE = (
    "# docker-compose.yml — production-like stack\n"
    "services:\n"
    "  zooplus-api:\n"
    "    build: .\n"
    "    ports: ['8080:8080']\n"
    "\n"
    "# Repo layout (FR5)\n"
    "cli/  src/  tests/  docs/  scripts/\n"
    "scripts/run_release_verify.ps1"
)

PANEL_NAMES = {
    "fr24": ("Panel FR24Code", "TextBox FR24Code"),
    "fr3rag": ("Panel FR3RagCode", "TextBox FR3RagCode"),
    "fr4": ("Panel FR4Code", "TextBox FR4Code"),
    "fr5": ("Panel FR5Code", "TextBox FR5Code"),
}


def _rgb(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _find_slide(prs: Presentation, *prefixes: str) -> int:
    for i, slide in enumerate(prs.slides):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            text = sh.text_frame.text.strip()
            if any(text.startswith(p) for p in prefixes):
                return i
    raise SystemExit(f"Slide not found: {prefixes!r}")


def _remove_named(slide, names: set[str]) -> None:
    for sh in list(slide.shapes):
        if sh.name in names:
            el = sh._element
            el.getparent().remove(el)


def _set_mono_block(text_frame, title: str, code: str) -> None:
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    h = text_frame.paragraphs[0]
    h.text = title
    h.font.name = "Arial"
    h.font.size = Pt(11)
    h.font.bold = True
    h.font.color.rgb = GREEN
    for line in code.split("\n"):
        p = text_frame.add_paragraph()
        p.text = line
        p.font.name = "Consolas"
        p.font.size = Pt(9)
        p.font.color.rgb = CODE_FG
        p.space_after = Pt(1)


def _upsert_code_panel(
    slide,
    *,
    key: str,
    title: str,
    code: str,
    left,
    top,
    width,
    height,
) -> None:
    panel_name, box_name = PANEL_NAMES[key]
    _remove_named(slide, {panel_name, box_name})
    panel = slide.shapes.add_shape(1, left, top, width, height)
    panel.name = panel_name
    _rgb(panel, CODE_BG)
    box = slide.shapes.add_textbox(
        left + Inches(0.12),
        top + Inches(0.1),
        width - Inches(0.24),
        height - Inches(0.15),
    )
    box.name = box_name
    _set_mono_block(box.text_frame, title, code)


def _shrink_picture(slide, *, top_height_inches: float = 3.05) -> None:
    for sh in slide.shapes:
        if sh.name == "Picture 14":
            sh.left = Inches(6.62)
            sh.top = Inches(1.45)
            sh.width = Inches(5.58)
            sh.height = Inches(top_height_inches)


def _resize_textbox_at_left(slide, left_emu: int, *, height_inches: float) -> None:
    for sh in slide.shapes:
        if sh.has_text_frame and sh.left == left_emu:
            sh.height = Inches(height_inches)


def apply_fr24_slide(prs: Presentation) -> None:
    slide = prs.slides[_find_slide(prs, "Response & RAG")]
    _shrink_picture(slide, top_height_inches=2.95)
    _upsert_code_panel(
        slide,
        key="fr24",
        title="FR2 + FR3 evidence (source)",
        code=FR24_CODE,
        left=Inches(6.62),
        top=Inches(4.55),
        width=Inches(5.58),
        height=Inches(2.2),
    )
    for sh in slide.shapes:
        if sh.has_text_frame and sh.name == "TextBox 12":
            text = sh.text or ""
            if "max 4 SKUs" in text:
                sh.text_frame.paragraphs[0].text = text.replace(
                    "max 4 SKUs per turn",
                    "default 4 picks · up to 20 if asked",
                )


def apply_fr3_rag_slide(prs: Presentation) -> None:
    slide = prs.slides[_find_slide(prs, "RAG strategy")]
    _resize_textbox_at_left(slide, 749_808, height_inches=3.75)
    _resize_textbox_at_left(slide, 5_989_320, height_inches=3.75)
    _upsert_code_panel(
        slide,
        key="fr3rag",
        title="FR3 — hybrid RAG + ingest (source)",
        code=FR3_RAG_CODE,
        left=Inches(6.55),
        top=Inches(5.05),
        width=Inches(5.9),
        height=Inches(1.55),
    )


def apply_fr4_slide(prs: Presentation) -> None:
    slide = prs.slides[_find_slide(prs, "Guardrails")]
    _shrink_picture(slide, top_height_inches=4.35)
    _upsert_code_panel(
        slide,
        key="fr4",
        title="FR4 — topic_check + constraints (source)",
        code=FR4_CODE,
        left=Inches(6.59),
        top=Inches(5.85),
        width=Inches(5.58),
        height=Inches(1.55),
    )


def apply_fr5_slide(prs: Presentation) -> None:
    slide = prs.slides[_find_slide(prs, "Production layout")]
    _shrink_picture(slide, top_height_inches=2.95)
    _upsert_code_panel(
        slide,
        key="fr5",
        title="FR5 — Docker + repo layout (source)",
        code=FR5_CODE,
        left=Inches(6.76),
        top=Inches(4.55),
        width=Inches(5.58),
        height=Inches(2.2),
    )


def apply_all_fr_code_panels(prs: Presentation) -> None:
    apply_fr24_slide(prs)
    apply_fr3_rag_slide(prs)
    apply_fr4_slide(prs)
    apply_fr5_slide(prs)


def main() -> None:
    if not DECK.is_file():
        raise SystemExit(f"Missing deck: {DECK}")
    prs = Presentation(str(DECK))
    apply_all_fr_code_panels(prs)
    tmp = DECK.with_name(f"{DECK.stem}_frcode{DECK.suffix}")
    prs.save(str(tmp))
    try:
        shutil.copy2(tmp, DECK)
        tmp.unlink(missing_ok=True)
        print(f"Added FR code panels to {DECK}")
    except PermissionError:
        print(f"PPT is open — saved to {tmp}")


if __name__ == "__main__":
    main()
