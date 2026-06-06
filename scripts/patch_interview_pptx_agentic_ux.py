#!/usr/bin/env python3
"""Update slides 6–9 — agentic UX content with clean layout (no overlapping text boxes)."""

from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

ROOT = Path(__file__).resolve().parents[1]
DECK = ROOT / "docs" / "deliverables" / "v0.1" / "zooplus-assistant-interview-15min-pro.pptx"

INK = RGBColor(0x33, 0x41, 0x55)
TITLE_INK = RGBColor(0x0A, 0x25, 0x40)

LEFT_PRIMARY_LEFT = 585_216
LEFT_DUPLICATE_LEFT = 658_368


def _find_slide(prs: Presentation, *title_prefixes: str) -> int:
    for i, slide in enumerate(prs.slides):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            text = sh.text.strip()
            if any(text.startswith(prefix) for prefix in title_prefixes):
                return i
    raise SystemExit(f"Slide not found: {title_prefixes!r}")


def _delete_shape(shape) -> None:
    el = shape._element
    el.getparent().remove(el)


def _remove_inner_left_column(slide) -> None:
    """Remove the overlapping inner bullet box (left=658368, mid-column only).

    Title, subtitle, and footer on the same left edge must stay.
    """
    to_remove = [
        sh
        for sh in slide.shapes
        if sh.has_text_frame
        and sh.left == LEFT_DUPLICATE_LEFT
        and sh.top
        and 1_000_000 < sh.top < 5_000_000
    ]
    for sh in to_remove:
        _delete_shape(sh)


def _primary_left_box(slide) -> object | None:
    for sh in slide.shapes:
        if sh.has_text_frame and sh.left == LEFT_PRIMARY_LEFT:
            if sh.top and 1_000_000 < sh.top < 5_000_000:
                return sh
    return None


def _set_paragraphs(text_frame, lines: list[tuple[str, bool]], *, size: int = 15) -> None:
    text_frame.clear()
    for i, (line, bold) in enumerate(lines):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = line
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = TITLE_INK if bold else INK
        p.space_after = Pt(4)
        if line.startswith("•"):
            p.level = 1


def _fill(shape, lines: list[tuple[str, bool]], *, size: int = 15) -> None:
    _set_paragraphs(shape.text_frame, lines, size=size)


def _fix_left_column(slide, lines: list[tuple[str, bool]], *, size: int = 15) -> None:
    _remove_inner_left_column(slide)
    primary = _primary_left_box(slide)
    if primary is not None:
        _fill(primary, lines, size=size)


def _right_column_bullets(slide) -> object | None:
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        if sh.left and sh.left > 5_500_000 and sh.top and 1_200_000 < sh.top < 6_000_000:
            if sh.height and sh.height > 2_000_000:
                return sh
    return None


SLIDE6_RIGHT = [
    ("Grounding & answer", True),
    ("• retrieved_products[] from same site_id hits only.", False),
    ("• Off-topic → empty list + polite decline (FR4).", False),
    ("• Synthesis: per-agent OpenCode LLM — never invent SKUs.", False),
    ("• Stream: /chat/stream → status → products → done.", False),
    ("• UI: one transient status bubble (agent shopper_status).", False),
    ("Production path", True),
    ("• routing_lexicon.json at ingest; optional Redis mirror.", False),
    ("• Managed vector DB + scheduled re-ingest (roadmap P1).", False),
]

SLIDE7_LEFT = [
    ("Agentic routing (v0.1.3)", True),
    ("• Conductor-first: classify TOPIC before catalog search.", False),
    ("• Social greeting → social agent only (no Chroma).", False),
    ("• Catalog lane → prefetch → process → synthesis.", False),
    ("• Catalog lexicon in prompt (brands/tokens from ingest).", False),
    ("Latency & cache", True),
    ("• TTL cache: intent · chat · retrieval.", False),
    ("• Optional Redis: shared cache + lexicon.", False),
    ("• Fast social: hola que tal in seconds, not 30s+.", False),
]

SLIDE8_LEFT = [
    ("OpenCode specialists", True),
    ("• conductor → lane + shopper_status (JSON).", False),
    ("• social-agent → multilingual social / decline.", False),
    ("• intent-agent → fallback cascade.", False),
    ("• rag + logic workers → hybrid retrieve, cap 4.", False),
    ("• synthesis → grounded prose from hits only.", False),
    ("Per-agent LLMs (opencode.json)", True),
    ("• Conductor / social: mimo-v2.5-free.", False),
    ("• Intent / RAG: deepseek-v4-flash.", False),
    ("• Logic: qwen3.6-plus · Topic-guard: minimax-m2.5.", False),
    ("• Synthesis: deepseek-v4-flash-free.", False),
    ("• UI badge shows real model from response meta.", False),
]

SLIDE9_LEFT = [
    ("FR4 guardrails + live demo", True),
    ("• Pet catalog only · default-deny off-topic.", False),
    ("• constraints.yaml + agentic intent + social decline.", False),
    ("Demo (http://127.0.0.1:8090/ui/)", True),
    ("• A: hola que tal → fast social (no RAG).", False),
    ("• B: cat food 40–60 EUR → stream status → products.", False),
    ("• C: weather / news → polite decline.", False),
    ("• Pick shop 1, 3, or 15 before sending.", False),
]


def main() -> None:
    if not DECK.is_file():
        raise SystemExit(f"Missing deck: {DECK}")
    prs = Presentation(str(DECK))

    i6 = _find_slide(prs, "RAG strategy")
    right6 = _right_column_bullets(prs.slides[i6])
    if right6 is not None:
        _fill(right6, SLIDE6_RIGHT, size=14)

    i7 = _find_slide(prs, "Agentic architecture", "Agentic routing")
    _fix_left_column(prs.slides[i7], SLIDE7_LEFT, size=14)
    for sh in prs.slides[i7].shapes:
        if sh.has_text_frame and "How each request" in (sh.text or ""):
            sh.text_frame.paragraphs[0].text = (
                "Conductor-first · catalog lexicon · stream status — fast multilingual UX"
            )

    i8 = _find_slide(prs, "Agents")
    _fix_left_column(prs.slides[i8], SLIDE8_LEFT, size=13)

    i9 = _find_slide(prs, "Guardrails", "FR4 guardrails")
    _fix_left_column(prs.slides[i9], SLIDE9_LEFT, size=14)

    tmp = DECK.with_name(f"{DECK.stem}_patched{DECK.suffix}")
    prs.save(str(tmp))
    try:
        shutil.copy2(tmp, DECK)
        tmp.unlink(missing_ok=True)
        print(f"Fixed layout + content in {DECK}")
    except PermissionError:
        print(f"PPT is open — saved to {tmp}")
        print("Close PowerPoint, then run: copy patched file over the deck.")


if __name__ == "__main__":
    main()
